import json
from pptx import Presentation
from PIL import Image
from pptx.util import Emu
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_slide_size(ppt):
    slide_width = ppt.slide_width.inches
    slide_height = ppt.slide_height.inches
    return slide_width, slide_height

def extract_master(ppt):
    master_data = []
    for master in ppt.slide_masters:
        background = master.background
        fill = background.fill
        fill_type = fill.type

        if fill_type == 1:
            if fill.fore_color.type == 0:
                fill_color = str(fill.fore_color.rgb)
            elif fill.fore_color.type == 1:
                fill_color = fill.fore_color.theme_color
            else:
                fill_color = None
        else:
            fill_color = None

        master_data.append({
            'name': master.name,
            'background_fill_type': fill_type,
            'background_fill_color': fill_color
        })
    return master_data

def extract_text(slide):
    text_data = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            for paragraph in text_frame.paragraphs:
                space_after = paragraph.space_after if paragraph.space_after else None
                for run in paragraph.runs:
                    font_size = run.font.size.pt if run.font.size else None
                    text_data.append({
                        'text': run.text,
                        'font_name': run.font.name,
                        'font_size': font_size,
                        'space_after': space_after,
                        'left': left,
                        'top': top,
                        'width': width,
                        'height': height
                    })
    return text_data


def extract_images(slide, slide_number):
    image_data = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture shape type
            image = shape.image
            image_path = f"slide_{slide_number}_image_{shape.shape_id}.png"
            with open(image_path, "wb") as f:
                f.write(image.blob)
            img = Image.open(image_path)
            width, height = img.size
            left, top = shape.left, shape.top
            crop_left, crop_right, crop_top, crop_bottom = shape.crop_left, shape.crop_right, shape.crop_top, shape.crop_bottom
            image_data.append({
                'path': image_path,
                'width': Emu(width * 9525),
                'height': Emu(height * 9525),
                'left': left,
                'top': top,
                'crop_left': crop_left,
                'crop_right': crop_right,
                'crop_top': crop_top,
                'crop_bottom': crop_bottom
            })
    return image_data


def pptPath_to_json(ppt_path):
    presentation_data = {}
    ppt = Presentation(ppt_path)

    slide_width, slide_height = extract_slide_size(ppt)
    presentation_data['slide_size'] = {'width': slide_width, 'height': slide_height}

    presentation_data['masters'] = extract_master(ppt)

    slides_data = []
    for slide_number, slide in enumerate(ppt.slides):
        slide_data = {
            'slide_number': slide_number + 1,
            'layout': slide.slide_layout.name,
            'text_data': extract_text(slide),
            'image_data': extract_images(slide, slide_number + 1)
        }
        slides_data.append(slide_data)

    presentation_data['slides'] = slides_data

    return json.dumps(presentation_data, indent=4)

def extract_elements(slide, slide_number):
    elements_data = []

    for index, shape in enumerate(slide.shapes):
        element = {'z_index': index}

        if shape.has_text_frame:
            text_frame = shape.text_frame

            if text_frame.paragraphs:
                paragraph = text_frame.paragraphs[0]
                if paragraph.runs:
                    run = paragraph.runs[0]
                    font_name = run.font.name
                    font_size = run.font.size
                else:
                    font_name = None
                    font_size = None
            else:
                font_name = None
                font_size = None

            element.update({
                'type': 'text',
                'text': text_frame.text,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'font_name': font_name,
                'font_size': font_size
            })

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element.update({
                'type': 'picture',
                'image_path': f'image_{slide_number}_{index}.png',
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            })

        elements_data.append(element)

    return elements_data


def ppt_to_json(file_path):
    ppt = pptx.Presentation(file_path)
    presentation_data = {
        'slide_width': ppt.slide_width,
        'slide_height': ppt.slide_height,
        'slides': []
    }

    for slide_number, slide in enumerate(ppt.slides, start=1):
        slide_data = {
            'slide_number': slide_number,
            'slide_layout': slide.slide_layout.name,
            'elements': extract_elements(slide, slide_number)
        }
        presentation_data['slides'].append(slide_data)

    return json.dumps(presentation_data, indent=4, default=str)



input_ppt = "test1.pptx"
output_json = "example.json"

with open(output_json, "w") as f:
    f.write(ppt_to_json(input_ppt))