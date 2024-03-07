import os
import csv
from pptx import Presentation
from PIL import Image, UnidentifiedImageError
import imghdr
import io
import logging
import imagehash
import pandas as pd


# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


def get_unique_hashes(csv_path):
    """使用pandas读取CSV文件，并返回所有'Is Duplicate'为0的图片hash的集合。"""
    if os.path.exists(csv_path):
        # 读取CSV文件
        df = pd.read_csv(csv_path)

        # 筛选出'Is Duplicate'为0的行
        unique_hashes_df = df[df['Is Duplicate'] == 0]

        # 如果DataFrame不为空，返回'Image Hash'列的值组成的集合
        if not unique_hashes_df.empty:
            return set(unique_hashes_df['Image Hash'])
    return set()


def mark_and_remove_duplicates_in_csv(file_path):
    """标记csv文件中hash重复的项目，重复标记为1，不重复标记为0"""
    # 读取CSV文件
    df = pd.read_csv(file_path)
    # 识别重复的hash，第一次出现标记为False（即1），重复出现标记为True（即0）
    df['Is Duplicate'] = df['Image Hash'].duplicated().astype(int)
    # 保存修改后的DataFrame到新的CSV文件，或者按需处理
    df.to_csv(file_path, index=False)

    # 找出所有标记为重复的文件
    duplicates = df[df['Is Duplicate'] == 1]['Image File']
    # 遍历并删除这些文件
    for file_path in duplicates:
        if os.path.exists(file_path):
            os.remove(file_path)
            print(f"Deleted: {file_path}")
        else:
            print(f"File not found, skipped: {file_path}")


def get_image_format(image_blob):
    """根据图片的二进制数据猜测图片格式。"""
    try:
        image_bytes = io.BytesIO(image_blob)
        img_format = imghdr.what(image_bytes)
        return img_format
    except Exception as e:
        logging.error(f"Error determining image format: {e}")
        return None


def is_size_similar(img_width, img_height, slide_width, slide_height):
    """判断图片尺寸是否符合要求。"""
    tolerance = 500
    return (abs(img_width - slide_width) < tolerance and abs(img_height - slide_height) < tolerance) or (img_width >= slide_width and img_height >= slide_height)


def save_image(image, slide_idx, shape_idx, pptx_filename, dest_folder, csv_writer, existing_hashes):
    """保存图片并记录到CSV。"""
    try:
        image_bytes = io.BytesIO(image.blob)
        img = Image.open(image_bytes)
        img_hash = str(imagehash.average_hash(img))

        img_format = get_image_format(image.blob) or 'png'
        file_extension = 'jpg' if img_format == 'jpeg' else img_format

        img_name = f"{os.path.splitext(os.path.basename(pptx_filename))[0]}_{slide_idx+1}_{shape_idx+1}.{file_extension}"
        img_path = os.path.join(dest_folder, img_name)

        if img_hash in existing_hashes:

            logging.info(f"重复图片: {img_path}")
            csv_writer.writerow([pptx_filename, img_path, img_hash])

        else:
            img.save(img_path)
            logging.info(f"保存图片: {img_path}")
            csv_writer.writerow([pptx_filename, img_path, img_hash])

    except UnidentifiedImageError:
        logging.error(f"UnidentifiedImageError: Cannot identify image file in {pptx_filename}, slide {slide_idx+1}, shape {shape_idx+1}.")
    except Exception as e:
        logging.error(f"Error saving image: {e}")


def save_slide_images(pptx_file, dest_folder, csv_writer, existing_hashes):
    """处理PPTX文件中的每个幻灯片图片。"""
    try:
        presentation = Presentation(pptx_file)
        for slide_idx, slide in enumerate(presentation.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # 图片类型
                    img_width = shape.width
                    img_height = shape.height
                    slide_width = presentation.slide_width
                    slide_height = presentation.slide_height

                    if is_size_similar(img_width, img_height, slide_width, slide_height):
                        save_image(shape.image, slide_idx, shape_idx, pptx_file, dest_folder, csv_writer, existing_hashes)
    except Exception as e:
        logging.error(f"Error processing {pptx_file}: {e}")


def main(src_folder, dest_folder, csv_file_path):
    """主函数，遍历目录，处理PPTX文件。"""
    if not os.path.exists(dest_folder):
        os.makedirs(dest_folder)
    try:
        existing_hashes = get_unique_hashes(csv_file_path)
        print(existing_hashes)
    except:
        existing_hashes = None

    with open(csv_file_path, mode='w', newline='', encoding='utf-8') as csv_file:
        csv_writer = csv.writer(csv_file)
        csv_writer.writerow(["PPTX File", "Image File", "Image Hash"])

        for root, dirs, files in os.walk(src_folder):
            for file in files:
                if file.endswith(".pptx") and not file.startswith("~$"):
                    pptx_path = os.path.join(root, file)
                    save_slide_images(pptx_path, dest_folder, csv_writer, existing_hashes)
                    logging.info(f"Processed {pptx_path}")

    # 找到重复的hash，更新csv表格，并删除重复的图片文件
    mark_and_remove_duplicates_in_csv(csv_file_path)


if __name__ == "__main__":
    src_folder = "/Users/birdmanoutman/上汽"
    dest_folder = "/Users/birdmanoutman/上汽/backgroundIMGsource"
    csv_file_path = os.path.join(dest_folder, "image_ppt_mapping.csv")
    main(src_folder, dest_folder, csv_file_path)
