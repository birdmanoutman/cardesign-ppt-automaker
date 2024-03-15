# findBackgroundIMG_V3.py
import io
import logging
import os
import sqlite3

import pandas as pd
from PIL import Image
from pptx import Presentation

# Ensure utils.py is in the same directory and contains calculate_hash, is_size_similar functions.
from utils import calculate_hash, is_size_similar

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


class SQLiteManager:
    def __init__(self, db_path):
        self.db_path = db_path
        self.conn = None
        self.cursor = None

    def connect(self):
        try:
            self.conn = sqlite3.connect(self.db_path)
            self.cursor = self.conn.cursor()
            logging.info("Database connection established.")
        except sqlite3.Error as e:
            logging.error(f"Error connecting to database: {e}")
            raise e

    def create_table(self, table_name):
        table_sql = f'''
            CREATE TABLE IF NOT EXISTS {table_name} (
                id INTEGER PRIMARY KEY,
                img_hash TEXT UNIQUE,
                img_path TEXT,
                pptx_path TEXT
            )'''
        try:
            self.cursor.execute(table_sql)
            self.conn.commit()
            logging.info("Table is ready.")
        except sqlite3.Error as e:
            logging.error(f"Error creating table: {e}")
            raise e

    def import_csv_to_database(self, csv_file_path, table_name):
        try:
            df = pd.read_csv(csv_file_path)
            for _, row in df.iterrows():
                img_hash = row['Image Hash']
                img_path = row['Image Path']
                pptx_path = row['PPTX Path']
                self.cursor.execute(
                    f'INSERT INTO {table_name} (img_hash, img_path, pptx_path) VALUES (?, ?, ?) ON CONFLICT(img_hash) DO REPLACE',
                    (img_hash, img_path, pptx_path))
            self.conn.commit()
            logging.info("CSV import completed.")
        except Exception as e:
            logging.error(f"Error importing CSV to database: {e}")
            raise e

    def check_duplicate(self, table_name, img_hash):
        try:
            self.cursor.execute(f"SELECT EXISTS(SELECT 1 FROM {table_name} WHERE img_hash=? LIMIT 1)", (img_hash,))
            return self.cursor.fetchone()[0]
        except sqlite3.Error as e:
            logging.error(f"Error checking duplicate: {e}")
            raise e

    def insert_image(self, table_name, img_hash, img_path, pptx_path):
        try:
            self.cursor.execute(
                f"INSERT INTO {table_name} (img_hash, img_path, pptx_path) VALUES (?, ?, ?) ON CONFLICT(img_hash) DO REPLACE",
                (img_hash, img_path, pptx_path))
            self.conn.commit()
        except sqlite3.Error as e:
            logging.error(f"Error inserting image: {e}")
            raise e

    def close(self):
        if self.cursor:
            self.cursor.close()
        if self.conn:
            self.conn.close()
        logging.info("Database connection closed.")


class ImageManager:
    def __init__(self, db_manager, table_name):
        self.db_manager = db_manager
        self.table_name = table_name

    def is_duplicate(self, img_hash):
        return self.db_manager.check_duplicate(self.table_name, img_hash)

    def save_image(self, img, img_hash, img_name, dest_folder, pptx_path):
        if self.is_duplicate(img_hash):
            logging.info(f"Duplicate image detected, not saved: {img_name}")
            return False
        else:
            if img.mode in ['RGBA', 'P']:
                img = img.convert('RGB')

            img_path = os.path.join(dest_folder, img_name)
            img.save(img_path)
            logging.info(f"Saved image: {img_path}")
            self.db_manager.insert_image(self.table_name, img_hash, img_path, pptx_path)
            return True


class ImageExtractor:
    def __init__(self, src_folder, dest_folder, image_manager):
        self.src_folder = src_folder
        self.dest_folder = dest_folder
        self.image_manager = image_manager

    def extract_images(self):
        for root, _, files in os.walk(self.src_folder):
            for file in files:
                if file.endswith(".pptx") and not file.startswith("~$"):
                    pptx_path = os.path.join(root, file)
                    self.process_pptx(pptx_path)

    def process_pptx(self, pptx_path):
        try:
            presentation = Presentation(pptx_path)
            for slide_idx, slide in enumerate(presentation.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == 13:  # Picture type
                        img = Image.open(io.BytesIO(shape.image.blob))
                        if is_size_similar(img, presentation):
                            img_hash = calculate_hash(shape.image.blob)
                            if img_hash:
                                img_name = f"{os.path.splitext(os.path.basename(pptx_path))[0]}_{slide_idx + 1}_{shape_idx + 1}.jpg"
                                self.image_manager.save_image(img, img_hash, img_name, self.dest_folder, pptx_path)
        except Exception as e:
            logging.error(f"Error processing PPTX: {e}")
            raise e


class ImageProcessor:
    def __init__(self, db_path, src_folder, dest_folder, csv_file_path, table_name):
        self.db_path = db_path
        self.src_folder = src_folder
        self.dest_folder = dest_folder
        self.csv_file_path = csv_file_path
        self.table_name = table_name

    def run(self):
        db_manager = SQLiteManager(self.db_path)
        db_manager.connect()
        db_manager.create_table(self.table_name)

        if self.csv_file_path:
            db_manager.import_csv_to_database(self.csv_file_path, self.table_name)

        image_manager = ImageManager(db_manager, self.table_name)

        if not os.path.exists(self.dest_folder):
            os.makedirs(self.dest_folder)

        extractor = ImageExtractor(self.src_folder, self.dest_folder, image_manager)
        extractor.extract_images()

        db_manager.close()


def main():
    # Adjust these parameters according to your setup
    src_folder = "your_source_folder"
    dest_folder = "your_destination_folder"
    db_path = "your_database_file.db"
    csv_file_path = "your_csv_file_path.csv"
    table_name = "your_table_name"

    image_processor = ImageProcessor(db_path, src_folder, dest_folder, csv_file_path, table_name)
    image_processor.run()


if __name__ == "__main__":
    main()
