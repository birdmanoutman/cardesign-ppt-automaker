"""这段代码是一个用于处理PPTX文件中的图片的Python脚本，名为findBackgroundIMG_V3.py。
SQLiteManager 负责数据库操作，ImageManager 负责图片的去重和保存逻辑，并与数据库交互。
ImageExtractor 负责从PPTX文件中提取图片，并通过 ImageManager 来处理这些图片。
ImageProcessor 作为协调者，初始化并连接了上述类的实例，定义了处理流程。
main 函数作为程序的启动点，配置并运行 ImageProcessor。
这种设计体现了关注点分离的原则，通过分离数据访问层、业务逻辑层和表示层，提高了代码的可维护性和可扩展性。"""

# findBackgroundIMG_sqliteV1.py
import io
import logging
import os
import sqlite3

import pandas as pd
from PIL import Image
from pptx import Presentation

# Ensure utils.py is in the same directory and contains calculate_hash, is_size_similar functions.
from utils import calculate_hash, is_size_similar

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


class SQLiteManager:
    """
    功能：管理SQLite数据库的连接、创建表、导入CSV数据到数据库、检查图片哈希值是否重复、插入图片记录到数据库，以及关闭数据库连接。
    与其他对象的关系：ImageManager 类会使用 SQLiteManager 类来查询数据库中的图片记录，判断图片是否重复，并存储新的图片记录。
    """
    def __init__(self, db_path):
        self.db_path = db_path
        self.conn = None
        self.cursor = None

    def connect(self):
        try:
            self.conn = sqlite3.connect(self.db_path)
            self.cursor = self.conn.cursor()
            logging.info("Database connection established.")
        except sqlite3.Error as e:
            logging.error(f"Error connecting to database: {e}")
            raise e

    def create_table(self, table_name):
        table_sql = f'''
            CREATE TABLE IF NOT EXISTS {table_name} (
                id INTEGER PRIMARY KEY,
                img_hash TEXT UNIQUE,
                img_path TEXT,
                pptx_path TEXT,
                is_duplicate INT
            )'''
        try:
            self.cursor.execute(table_sql)
            self.conn.commit()
            logging.info("Table is ready.")
        except sqlite3.Error as e:
            logging.error(f"Error creating table: {e}")
            raise e

    def import_csv_to_database(self, csv_file_path, table_name):
        try:
            df = pd.read_csv(csv_file_path)
            for _, row in df.iterrows():
                img_hash = row['Image Hash']
                img_path = row['Image File']
                pptx_path = row['PPTX File']
                is_duplicate = row['Is Duplicate']
                self.cursor.execute(
                    f'INSERT INTO {table_name} (img_hash, img_path, pptx_path, is_duplicate) VALUES (?, ?, ?,?'
                    f') ON CONFLICT('
                    f'img_hash) DO UPDATE SET img_path=excluded.img_path, pptx_path=excluded.pptx_path',
                    (img_hash, img_path, pptx_path, is_duplicate))

            self.conn.commit()
            logging.info("CSV import completed.")
        except Exception as e:
            logging.error(f"Error importing CSV to database: {e}")
            raise e

    def check_duplicate(self, table_name, img_hash):
        try:
            self.cursor.execute(f"SELECT EXISTS(SELECT 1 FROM {table_name} WHERE img_hash=? LIMIT 1)", (img_hash,))
            return self.cursor.fetchone()[0]
        except sqlite3.Error as e:
            logging.error(f"Error checking duplicate: {e}")
            raise e

    def insert_image(self, table_name, img_hash, img_path, pptx_path):
        try:
            self.cursor.execute(
                f"INSERT INTO {table_name} (img_hash, img_path, pptx_path) VALUES (?, ?, ?) ON CONFLICT(img_hash) DO REPLACE",
                (img_hash, img_path, pptx_path))
            self.conn.commit()
        except sqlite3.Error as e:
            logging.error(f"Error inserting image: {e}")
            raise e

    def close(self):
        if self.cursor:
            self.cursor.close()
        if self.conn:
            self.conn.close()
        logging.info("Database connection closed.")


class ImageManager:
    """
    功能：负责检查图片是否为重复图片（通过图片的哈希值），保存新的图片到指定文件夹，并更新数据库中的记录。
    与其他对象的关系：依赖 SQLiteManager 类来执行数据库操作，由 ImageExtractor 类调用来处理具体的图片保存逻辑。
    """

    def __init__(self, db_manager, table_name):
        self.db_manager = db_manager
        self.table_name = table_name

    def is_duplicate(self, img_hash):
        return self.db_manager.check_duplicate(self.table_name, img_hash)

    def save_image(self, img, img_hash, img_name, dest_folder, pptx_path):
        if self.is_duplicate(img_hash):
            logging.info(f"Duplicate image detected, not saved: {img_name}")
            return False
        else:
            if img.mode in ['RGBA', 'P']:
                img = img.convert('RGB')

            img_path = os.path.join(dest_folder, img_name)
            img.save(img_path)
            logging.info(f"Saved image: {img_path}")
            self.db_manager.insert_image(self.table_name, img_hash, img_path, pptx_path)
            return True


class ImageExtractor:
    """
    功能：遍历指定源文件夹中的所有PPTX文件，提取每个演示文稿中的图片，并处理每张图片。
    与其他对象的关系：使用 ImageManager 类来处理图片的保存和数据库记录更新。它是图片提取过程的起点，负责具体的图片提取逻辑。
    """

    def __init__(self, src_folder, dest_folder, image_manager):
        self.src_folder = src_folder
        self.dest_folder = dest_folder
        self.image_manager = image_manager

    def extract_images(self):
        for root, _, files in os.walk(self.src_folder):
            for file in files:
                if file.endswith(".pptx") and not file.startswith("~$"):
                    pptx_path = os.path.join(root, file)
                    self.process_pptx(pptx_path)

    def process_pptx(self, pptx_path):
        try:
            presentation = Presentation(pptx_path)
            for slide_idx, slide in enumerate(presentation.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == 13:  # Picture type
                        img = Image.open(io.BytesIO(shape.image.blob))
                        if is_size_similar(img, presentation):
                            img_hash = calculate_hash(shape.image.blob)
                            if img_hash:
                                img_name = f"{os.path.splitext(os.path.basename(pptx_path))[0]}_{slide_idx + 1}_{shape_idx + 1}.jpg"
                                self.image_manager.save_image(img, img_hash, img_name, self.dest_folder, pptx_path)
        except Exception as e:
            logging.error(f"Error processing PPTX: {e}")
            raise e


class ImageProcessor:
    """
    功能：这是整个脚本的核心类，负责初始化数据库管理器、图片管理器和图片提取器。它设置了整个图片处理流程的参数，如源文件夹、目标文件夹、数据库路径、CSV文件路径和数据库表名，并触发图片提取和处理过程。
    与其他对象的关系：它实例化了 SQLiteManager, ImageManager, 和 ImageExtractor 类，并通过它们协同工作完成整个图片处理流程。
    """

    def __init__(self, db_path, src_folder, dest_folder, csv_file_path, table_name):
        self.db_path = db_path
        self.src_folder = src_folder
        self.dest_folder = dest_folder
        self.csv_file_path = csv_file_path
        self.table_name = table_name

    def run(self):
        db_manager = SQLiteManager(self.db_path)
        db_manager.connect()
        db_manager.create_table(self.table_name)

        if self.csv_file_path:
            db_manager.import_csv_to_database(self.csv_file_path, self.table_name)

        image_manager = ImageManager(db_manager, self.table_name)

        if not os.path.exists(self.dest_folder):
            os.makedirs(self.dest_folder)

        extractor = ImageExtractor(self.src_folder, self.dest_folder, image_manager)
        extractor.extract_images()

        db_manager.close()


def main():
    # 已知参数
    db_path = "image_gallery.db"  # 数据库文件路径
    table_name = "image_ppt_mapping"  # 数据库表名
    csv_file_path = r"D:\PycharmProjects\cardesign-ppt-automaker\findBackgroundIMG\source\image_ppt_mapping.csv"  #
    # CSV文件路径

    # 首次迁移CSV数据
    # 1. 初始化数据库管理器并连接数据库
    db_manager = SQLiteManager(db_path)
    db_manager.connect()
    # 2. 创建表（如果尚不存在）
    db_manager.create_table(table_name)
    # 3. 从CSV文件导入数据
    db_manager.import_csv_to_database(csv_file_path, table_name)


if __name__ == "__main__":
    main()
