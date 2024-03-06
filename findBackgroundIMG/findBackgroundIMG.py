import os
from pptx import Presentation
from PIL import Image
import imghdr
import io
import csv

def get_image_format(image_blob):
    """
    根据图片的二进制数据猜测图片格式。
    返回值是一个字符串，如 'jpeg', 'png' 等。
    """
    # 将图片二进制数据转换为字节流
    image_bytes = io.BytesIO(image_blob)
    # 使用imghdr检测图片格式
    img_format = imghdr.what(image_bytes)
    return img_format


def is_size_similar(img_width, img_height, slide_width, slide_height):
    tolerance = 500  # 容忍度，可根据需要调整
    if abs(img_width - slide_width) < tolerance and abs(img_height - slide_height) < tolerance:
        return True
    elif img_width >= slide_width and img_height >= slide_height:
        return True
    return False

def save_image(image, slide_idx, shape_idx, pptx_filename, dest_folder, csv_writer):
    image_bytes = io.BytesIO(image.blob)
    img = Image.open(image_bytes)

    # 猜测图片格式
    img_format = get_image_format(image.blob)
    # 如果无法确定格式，默认使用JPEG
    if img_format not in ['jpeg', 'png']:
        img_format = 'png'
    # 根据格式设置文件后缀
    file_extension = 'jpg' if img_format == 'jpeg' else img_format

    img_name = f"{os.path.splitext(os.path.basename(pptx_filename))[0]}_{slide_idx+1}_{shape_idx+1}.{file_extension}"
    img_path = os.path.join(dest_folder, img_name)

    # 保存图片，使用推断出的格式
    img.save(img_path)
    print(f"图片已保存: {img_path}")

    # 写入图片保存路径和对应的PPTX文件路径到CSV
    csv_writer.writerow([pptx_filename, img_path])

def save_slide_images(pptx_file, dest_folder, csv_writer):
    presentation = Presentation(pptx_file)
    for slide_idx, slide in enumerate(presentation.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # 图片类型
                img_width = shape.width
                img_height = shape.height
                slide_width = presentation.slide_width
                slide_height = presentation.slide_height

                if is_size_similar(img_width, img_height, slide_width, slide_height):
                    save_image(shape.image, slide_idx, shape_idx, pptx_file, dest_folder, csv_writer)

# 指定需要遍历的文件夹路径和存储图片的文件夹路径
src_folder = "/Users/birdmanoutman/Desktop"
dest_folder = "/Users/birdmanoutman/Desktop/outputTest"
csv_file_path = os.path.join(dest_folder, "image_ppt_mapping.csv")


# 确保目标文件夹存在
if not os.path.exists(dest_folder):
    os.makedirs(dest_folder)

with open(csv_file_path, mode='w', newline='', encoding='utf-8') as csv_file:
    csv_writer = csv.writer(csv_file)
    csv_writer.writerow(["PPTX File", "Image File"])  # 写入表头

    for root, dirs, files in os.walk(src_folder):
        for file in files:
            if file.endswith(".pptx") and not file.startswith("~$"):
                pptx_path = os.path.join(root, file)
                save_slide_images(pptx_path, dest_folder, csv_writer)