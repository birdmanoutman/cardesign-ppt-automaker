# findBackgroundIMGV3.py
import io
from PIL import Image
import imagehash
import logging
from pptx import Presentation
import os

import pandas as pd
import csv


class ImageManager:
    def __init__(self, existing_hashes):
        self.existing_hashes = existing_hashes

    def calculate_hash(self, image_blob):
        try:
            image_bytes = io.BytesIO(image_blob)
            img = Image.open(image_bytes)
            return str(imagehash.average_hash(img))
        except Exception as e:
            logging.error(f"Error calculating image hash: {e}")
            return None

    def is_duplicate(self, img_hash):
        return img_hash in self.existing_hashes

    def save_image(self, img, img_hash, img_name, dest_folder):
        if self.is_duplicate(img_hash):
            logging.info(f"Duplicate image detected, not saved: {img_name}")
            return False
        else:
            # 如果图片模式是'RGBA'或'P'，转换为'RGB'
            if img.mode == 'RGBA' or img.mode == 'P':
                img = img.convert('RGB')

            img_path = os.path.join(dest_folder, img_name)
            img.save(img_path)
            logging.info(f"Saved image: {img_path}")
            self.existing_hashes.add(img_hash)
            return True


class CSVManager:
    def __init__(self, csv_file_path):
        self.csv_file_path = csv_file_path

    def get_existing_hashes(self):
        if not os.path.exists(self.csv_file_path):
            return set()
        df = pd.read_csv(self.csv_file_path)
        return set(df['Image Hash'])

    def write_to_csv(self, row):
        with open(self.csv_file_path, mode='a', newline='', encoding='utf-8') as csv_file:
            csv_writer = csv.writer(csv_file)
            csv_writer.writerow(row)


class ImageExtractor:
    def __init__(self, src_folder, dest_folder, image_manager, csv_manager):
        self.src_folder = src_folder
        self.dest_folder = dest_folder
        self.image_manager = image_manager
        self.csv_manager = csv_manager

    def extract_images(self):
        for root, _, files in os.walk(self.src_folder):
            for file in files:
                if file.endswith(".pptx") and not file.startswith("~$"):
                    pptx_path = os.path.join(root, file)
                    self.process_pptx(pptx_path)

    def process_pptx(self, pptx_path):
        presentation = Presentation(pptx_path)
        for slide_idx, slide in enumerate(presentation.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # 图片类型
                    # 获取图片和幻灯片的尺寸
                    img_width = shape.width
                    img_height = shape.height
                    slide_width = presentation.slide_width
                    slide_height = presentation.slide_height

                    # 只处理符合尺寸要求的图片
                    if self.is_size_similar(img_width, img_height, slide_width, slide_height):
                        self.process_image(shape.image, slide_idx, shape_idx, pptx_path)


    def process_image(self, image, slide_idx, shape_idx, pptx_path):
        img_hash = self.image_manager.calculate_hash(image.blob)
        if img_hash:
            img_name = f"{os.path.splitext(os.path.basename(pptx_path))[0]}_{slide_idx+1}_{shape_idx+1}.jpg"
            saved = self.image_manager.save_image(Image.open(io.BytesIO(image.blob)), img_hash, img_name, self.dest_folder)
            self.csv_manager.write_to_csv([pptx_path, os.path.join(self.dest_folder, img_name) if saved else 'Duplicate', img_hash])

    def is_size_similar(self, img_width, img_height, slide_width, slide_height):
        """判断图片尺寸是否符合要求。"""
        tolerance = 500
        return (abs(img_width - slide_width) < tolerance and abs(img_height - slide_height) < tolerance) or (img_width >= slide_width and img_height >= slide_height)


def main():
    src_folder = "/Volumes/Backup/mac_backup"
    dest_folder = "/Users/birdmanoutman/上汽/backgroundIMGsource"
    csv_file_path = os.path.join(dest_folder, "image_ppt_mapping.csv")

    csv_manager = CSVManager(csv_file_path)
    existing_hashes = csv_manager.get_existing_hashes()
    image_manager = ImageManager(existing_hashes)

    if not os.path.exists(dest_folder):
        os.makedirs(dest_folder)
    if not os.path.exists(csv_file_path):
        csv_manager.write_to_csv(["PPTX File", "Image File", "Image Hash"])

    extractor = ImageExtractor(src_folder, dest_folder, image_manager, csv_manager)
    extractor.extract_images()


if __name__ == "__main__":
    main()
