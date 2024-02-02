from pptx import Presentation
from io import BytesIO
import base64
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
import cleanppt

inputPath = 'test.pptx'
outputPath = 'test_output.pptx'


def get_presentation_size(pptx_path):
    ppt = Presentation(pptx_path)
    width = ppt.slide_width
    height = ppt.slide_height
    return width, height


def create_presentation(output_path, width, height):
    # 创建一个新的 PowerPoint 文档
    ppt = Presentation()
    # 设置文档的尺寸
    ppt.slide_width = width
    ppt.slide_height = height
    ppt.save(output_path)
    return ppt


def get_slide_image_info(input_path):
    """输入文件路径
    返回文件第一页所有图片的信息"""
    ppt = Presentation(input_path)
    imgDatalist = []
    print(len(ppt.slides[0].shapes))
    for shape in ppt.slides[0].shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # 检查尺寸和位置是否有效
            if shape.left and shape.top and shape.width and shape.height:
                original_left = shape.left
                original_top = shape.top
                original_width = shape.width
                original_height = shape.height
                placeInfo = [original_left, original_top, original_width, original_height]

                # 获取原始图片裁剪信息（如果有的话）
                crop_left = shape.crop_left
                crop_top = shape.crop_top
                crop_right = shape.crop_right
                crop_bottom = shape.crop_bottom
                cropInfo = [crop_left, crop_top, crop_right, crop_bottom]

                # 获取图像的Base64编码
                image_data = shape.image.blob
                image_base64 = base64.b64encode(image_data).decode('utf-8')

                # 将Base64编码转换为BytesIO对象
                image_bytes = BytesIO(base64.b64decode(image_base64))
                imgDatalist.append([image_bytes, placeInfo, cropInfo])
    return imgDatalist


def place_img(ppt, slideNumber, image_bytes, placeInfo, cropInfo):
    slide = ppt.slides[slideNumber]

    # 复制图片并调整位置
    new_shape = slide.shapes.add_picture(image_bytes, *placeInfo)

    # 调整新图片的裁剪信息
    crop_left = cropInfo[0]
    crop_top = cropInfo[1]
    crop_right = cropInfo[2]
    crop_bottom = cropInfo[3]

    if crop_left or crop_top or crop_right or crop_bottom:
        new_shape.crop_left = crop_left
        new_shape.crop_top = crop_top
        new_shape.crop_right = crop_right
        new_shape.crop_bottom = crop_bottom
        print('cropped')

    ppt.save(outputPath)


width, height = get_presentation_size(inputPath)
print(width, height)

# ppt2 = create_presentation(outputPath, width, height)
# slide_layout = ppt2.slide_layouts[6]
# ppt2.slides.add_slide(slide_layout)
# ppt2.save(outputPath)




