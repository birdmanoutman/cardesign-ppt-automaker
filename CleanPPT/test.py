from pptx import Presentation


def cleanGarbageLayouts(path):
    # 打开现有的PPTX文件
    ppt = Presentation(path)

    # 用于存储不同布局名称的集合（避免重复）
    layout_names = set()

    # 遍历所有幻灯片布局并记录布局名称
    for slide_layout in ppt.slide_master.slide_layouts:
        layout_names.add(slide_layout.name)

    # 将布局名称转换为列表并排序
    layout_names_list = sorted(list(layout_names))

    layoutNeedsToDeleted = []
    for exist_layout in layout_names_list:
        layout_used = False
        for slide in ppt.slides:
            if slide.slide_layout.name == exist_layout:
                layout_used = True
                for slide_layout in ppt.slide_master.slide_layouts:
                    if slide_layout.name == exist_layout:
                        # 遍历模板中的所有占位符并删除它们
                        for placeholder in slide_layout.placeholders:
                            sp = placeholder._element
                            sp.getparent().remove(sp)
                break
        if not layout_used:
            layoutNeedsToDeleted.append(exist_layout)
    # 遍历所有幻灯片布局并删除它们
    for slide_layout in ppt.slide_layouts:
        if slide_layout.name in layoutNeedsToDeleted:  # 保留 "Blank" 布局
            ppt.slide_master.slide_layouts.remove(slide_layout)
    ppt.save(path)


cleanGarbageLayouts('//introducing_rose_flowers.pptx')


