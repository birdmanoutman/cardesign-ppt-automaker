import json
import os
from pptx import Presentation
from PIL import Image as PILImage
import io
from wand.image import Image as WandImage
from pptx.enum.dml import MSO_FILL_TYPE, MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.graphfrm import GraphicFrame



def convert_wmf_to_png(wmf_path, output_path):
    with WandImage(filename=wmf_path) as img:
        img.format = 'png'
        img.save(filename=output_path)


def save_image(image, slide_index, image_index, output_dir="images"):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    image_format = image.format
    # 假定WMF图像以二进制形式存在
    if image_format == 'WMF':
        wmf_path = os.path.join(output_dir, f"temp_slide_{slide_index}_image_{image_index}.wmf")
        with open(wmf_path, "wb") as wmf_file:
            wmf_file.write(image.blob)
        # 转换WMF到PNG
        png_path = os.path.join(output_dir, f"slide_{slide_index}_image_{image_index}.png")
        convert_wmf_to_png(wmf_path, png_path)
        # 清理临时WMF文件
        os.remove(wmf_path)
        return png_path
    else:
        # 为图片构建文件路径，这里假设保存为PNG格式
        image_path = os.path.join(output_dir, f"slide_{slide_index}_image_{image_index}.png")
        # 使用PIL的save方法保存图像
        image.save(image_path, format='PNG')
        return image_path


def extract_text(shape, shape_index):
    """
    从PPT元素中提取文本及其属性。
    """
    text_details = {
        "type": "text",
        "z_index": shape_index,
        "content": shape.text,
        "left": shape.left.pt,
        "top": shape.top.pt,
        "width": shape.width.pt,
        "height": shape.height.pt,
        # 可以添加更多文本特定的属性，如字体大小、颜色等
    }

    if shape.has_text_frame and shape.text_frame.paragraphs:
        first_paragraph = shape.text_frame.paragraphs[0]
        if first_paragraph.runs:
            first_run = first_paragraph.runs[0]
            text_details["font_name"] = first_run.font.name
            text_details["font_size"] = first_run.font.size.pt if first_run.font.size else None
            text_details["bold"] = first_run.font.bold
            text_details["italic"] = first_run.font.italic
            text_details["underline"] = first_run.font.underline

            # 检查并安全地访问字体颜色的.rgb属性
            font_color = first_run.font.color
            if font_color.type == 'RGB' or font_color.type == 'SCHEME':
                text_details["color"] = font_color.rgb
            else:
                text_details["color"] = None  # 或设置为默认颜色

    return text_details


def extract_shape(shape, shape_index):
    shape_details = {
        "type": "shape",
        "z_index": shape_index,
        "left": shape.left.pt,
        "top": shape.top.pt,
        "width": shape.width.pt,
        "height": shape.height.pt,
    }

    # 检查形状是否支持填充属性，并记录填充颜色
    if hasattr(shape, 'fill') and shape.fill.type == MSO_FILL_TYPE.SOLID:
        shape_details['fill_color'] = shape.fill.fore_color.rgb if shape.fill.fore_color.type == MSO_COLOR_TYPE.RGB else None

    # 检查形状是否有线条属性，并记录边框颜色和粗细
    if hasattr(shape, 'line') and shape.line.color.type == MSO_COLOR_TYPE.RGB:
        shape_details['line_color'] = shape.line.color.rgb
        shape_details['line_width'] = shape.line.width.pt

    # 记录阴影效果（如果有）
    # 仅在形状支持阴影属性时尝试访问
    if not isinstance(shape, GraphicFrame) and hasattr(shape, 'shadow'):
        try:
            if not shape.shadow.inherit:
                shape_details['shadow'] = {
                    "effect": shape.shadow.effect,
                    "angle": shape.shadow.angle,
                    "distance": shape.shadow.distance.pt,
                }
        except NotImplementedError:
            # 对于不支持阴影属性的形状，可以在这里处理异常
            pass


    # 根据形状类型处理特定属性
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        shape_details['auto_shape_type'] = shape.auto_shape_type
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        # 特殊处理图片类型
        shape_details['picture'] = "Special handling for pictures might be required here."
    # 为其他特定形状类型添加更多的条件分支

    return shape_details


def extract_images(slide, slide_index, output_dir):
    images = []
    for shape_index, shape in enumerate(slide.shapes):
        if shape.shape_type == 13:  # 图片类型
            image_bytes = io.BytesIO(shape.image.blob)
            try:
                with PILImage.open(image_bytes) as img:
                    # 计算图片的位置和尺寸
                    left = shape.left.pt
                    top = shape.top.pt
                    width = shape.width.pt
                    height = shape.height.pt
                    print('图片宽：{},高:{}'.format(width, height))

                    if img.format == 'WMF':
                        # 创建临时 WMF 文件
                        wmf_path = os.path.join(output_dir, f"temp_slide_{slide_index}_image_{shape_index}.wmf")
                        with open(wmf_path, 'wb') as f:
                            f.write(shape.image.blob)
                        # 转换 WMF 到 PNG
                        png_path = os.path.join(output_dir, f"slide_{slide_index}_image_{shape_index}.png")
                        convert_wmf_to_png(wmf_path, png_path)
                        image_path = png_path
                        # 清理临时 WMF 文件
                        os.remove(wmf_path)
                    else:
                        # 直接保存非 WMF 图片
                        image_path = os.path.join(output_dir, f"slide_{slide_index}_image_{shape_index}.png")
                        img.save(image_path, format='PNG')

                    images.append({
                        "type": "image",
                        "z_index": shape_index,  # 添加z_index属性
                        "path": image_path,
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height
                    })
            except Exception as e:
                print(f"Error processing image on slide {slide_index}, shape {shape_index}: {e}")
    return images


# 确保之前定义的extract_images和extract_shape函数适当地处理了图片和形状
def pptx_to_json(pptx_path, json_path, images_dir="images"):
    presentation = Presentation(pptx_path)
    presentation_data = {
        "presentation": {
            "slides": [],
            "slide_size": {"width": presentation.slide_width.pt, "height": presentation.slide_height.pt}
        }
    }

    for slide_index, slide in enumerate(presentation.slides):
        slide_data = {
            "slide_number": slide_index + 1,
            "elements": []
        }

        # 直接为每个幻灯片调用extract_images
        images = extract_images(slide, slide_index, images_dir)
        slide_data["elements"].extend(images)  # 将图片添加到元素列表中

        for shape_index, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text_details = extract_text(shape, shape_index)
                text_details["z_index"] = shape_index  # 添加 z_index
                slide_data["elements"].append(text_details)
            elif not shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # 排除图片类型，因为它们已被处理
                shape_info = extract_shape(shape, shape_index)
                shape_info["z_index"] = shape_index  # 添加 z_index
                slide_data["elements"].append(shape_info)

        presentation_data["presentation"]["slides"].append(slide_data)

    with open(json_path, "w") as json_file:
        json.dump(presentation_data, json_file, ensure_ascii=False, indent=4)


# 使用示例
pptx_path = "/Users/birdmanoutman/Desktop/项目/门板空间集成/20240125 空间集成_门板集成设计.pptx"
json_path = "output_presentation.json"
pptx_to_json(pptx_path, json_path)
