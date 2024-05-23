#!/usr/bin/env python
from pptx import Presentation
import sys


def main():
    if len(sys.argv) != 2:
        print("Usage: {} <file_path>".format(sys.argv[0]))
        sys.exit(1)

    file_path = sys.argv[1]
    cleanGarbageLayouts(file_path)

    # 在这里执行您的脚本逻辑，使用file_path参数


def cleanGarbageLayouts(path):
    """输入：ppt文件路径，
    输出：把原来的文件中没用的模板全部清除掉，同时删除所有有用模板中的占位符"""
    ppt = Presentation(path)
    # 删除主模板（第一个幻灯片母板）中的所有占位符
    slide_master = ppt.slide_master

    # 遍历主幻灯片母板中的占位符并删除它们
    for placeholder in slide_master.placeholders:
        sp = placeholder._element
        sp.getparent().remove(sp)

    # 收集所有幻灯片布局的名称
    all_layouts = set(layout.name for layout in ppt.slide_master.slide_layouts)

    # 检查哪些布局没有被使用
    unused_layouts = set(all_layouts)
    for slide in ppt.slides:
        unused_layouts.discard(slide.slide_layout.name)

    # 删除未使用的布局
    for layout_name in unused_layouts:
        ppt.slide_master.slide_layouts.remove(ppt.slide_master.slide_layouts.get_by_name(layout_name))

    # 删除已经被使用的布局中的占位符
    for slide in ppt.slides:
        slide_layout = ppt.slide_master.slide_layouts.get_by_name(slide.slide_layout.name)
        for placeholder in slide_layout.placeholders:
            sp = placeholder._element
            sp.getparent().remove(sp)
    ppt.save(path)


# 调用函数并传递PPT文件路径
# cleanGarbageLayouts('/Users/birdmanoutman/PycharmProjects/cardesign-ppt-automaker/introducing_rose_flowers.pptx')


if __name__ == "__main__":
    main()


