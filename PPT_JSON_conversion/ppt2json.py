import json
import os
from pptx import Presentation
from PIL import Image as PILImage
import io
from wand.image import Image as WandImage
from pptx.util import Inches


def convert_wmf_to_png(wmf_path, output_path):
    with WandImage(filename=wmf_path) as img:
        img.format = 'png'
        img.save(filename=output_path)


def save_image(image, slide_index, image_index, output_dir="images"):
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    image_format = image.format
    # 假定WMF图像以二进制形式存在
    if image_format == 'WMF':
        wmf_path = os.path.join(output_dir, f"temp_slide_{slide_index}_image_{image_index}.wmf")
        with open(wmf_path, "wb") as wmf_file:
            wmf_file.write(image.blob)
        # 转换WMF到PNG
        png_path = os.path.join(output_dir, f"slide_{slide_index}_image_{image_index}.png")
        convert_wmf_to_png(wmf_path, png_path)
        # 清理临时WMF文件
        os.remove(wmf_path)
        return png_path
    else:
        # 为图片构建文件路径，这里假设保存为PNG格式
        image_path = os.path.join(output_dir, f"slide_{slide_index}_image_{image_index}.png")
        # 使用PIL的save方法保存图像
        image.save(image_path, format='PNG')
        return image_path


def extract_text(shape):
    """
    从PPT元素中提取文本及其属性。
    """
    text_details = {
        "type": "text",
        "content": shape.text,
        "left": shape.left.pt,
        "top": shape.top.pt,
        "width": shape.width.pt,
        "height": shape.height.pt
    }

    if shape.has_text_frame and shape.text_frame.paragraphs:
        first_paragraph = shape.text_frame.paragraphs[0]
        if first_paragraph.runs:
            first_run = first_paragraph.runs[0]
            text_details["font_name"] = first_run.font.name
            text_details["font_size"] = first_run.font.size.pt if first_run.font.size else None
            text_details["bold"] = first_run.font.bold
            text_details["italic"] = first_run.font.italic
            text_details["underline"] = first_run.font.underline

            # 检查并安全地访问字体颜色的.rgb属性
            font_color = first_run.font.color
            if font_color.type == 'RGB' or font_color.type == 'SCHEME':
                text_details["color"] = font_color.rgb
            else:
                text_details["color"] = None  # 或设置为默认颜色

    return text_details


def extract_shape(element):
    """
    提取形状信息。
    """
    # 这里可以根据需要扩展，提取形状的详细信息
    shape_info = {
        "type": "shape",
        "shape_type": element.shape_type
    }
    return shape_info
def extract_images(slide, slide_index, output_dir):
    images = []
    for shape_index, shape in enumerate(slide.shapes):
        if shape.shape_type == 13:  # 图片类型
            image_bytes = io.BytesIO(shape.image.blob)
            try:
                with PILImage.open(image_bytes) as img:
                    # 计算图片的位置和尺寸
                    left = Inches(shape.left).inches
                    top = Inches(shape.top).inches
                    width = Inches(shape.width).inches
                    height = Inches(shape.height).inches

                    if img.format == 'WMF':
                        # 创建临时 WMF 文件
                        wmf_path = os.path.join(output_dir, f"temp_slide_{slide_index}_image_{shape_index}.wmf")
                        with open(wmf_path, 'wb') as f:
                            f.write(shape.image.blob)
                        # 转换 WMF 到 PNG
                        png_path = os.path.join(output_dir, f"slide_{slide_index}_image_{shape_index}.png")
                        convert_wmf_to_png(wmf_path, png_path)
                        image_path = png_path
                        # 清理临时 WMF 文件
                        os.remove(wmf_path)
                    else:
                        # 直接保存非 WMF 图片
                        image_path = os.path.join(output_dir, f"slide_{slide_index}_image_{shape_index}.png")
                        img.save(image_path, format='PNG')

                    images.append({
                        "type": "image",
                        "path": image_path,
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height
                    })
            except Exception as e:
                print(f"Error processing image on slide {slide_index}, shape {shape_index}: {e}")
    return images


# 确保之前定义的extract_images和extract_shape函数适当地处理了图片和形状

def pptx_to_json(pptx_path, json_path, images_dir="images"):
    """
    将PPTX文件转换为JSON格式，包括过渡效果、动画和详细的形状信息等。
    """
    presentation = Presentation(pptx_path)
    presentation_data = {
        "presentation": {
            "slides": [],
            "slide_size": {"width": presentation.slide_width.pt, "height": presentation.slide_height.pt}
        }
    }

    for slide_index, slide in enumerate(presentation.slides):
        slide_data = {
            "slide_number": slide_index + 1,
            "elements": []
        }

        # 处理图片和形状
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 图片类型
                slide_data["elements"].extend(extract_images(slide, slide_index, images_dir))
            elif shape.has_text_frame:
                text_details = extract_text(shape)
                slide_data["elements"].append(text_details)
            else:  # 其他形状类型，如纯形状等
                shape_info = extract_shape(shape)  # 确保这个函数能够提取形状的详细信息
                slide_data["elements"].append(shape_info)

        presentation_data["presentation"]["slides"].append(slide_data)

    with open(json_path, "w") as json_file:
        json.dump(presentation_data, json_file, ensure_ascii=False, indent=4)


# 使用示例
pptx_path = "/Users/birdmanoutman/Desktop/项目/门板空间集成/20240125 空间集成_门板集成设计.pptx"
json_path = "output_presentation.json"
pptx_to_json(pptx_path, json_path)
