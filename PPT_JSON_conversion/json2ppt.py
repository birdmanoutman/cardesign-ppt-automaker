import json
from pptx import Presentation
from pptx.util import Pt  # 使用点数作为单位
from pptx.dml.color import RGBColor  # 导入 RGBColor

def add_text(slide, element):
    left = Pt(element['left'])  # 使用点数
    top = Pt(element['top'])  # 使用点数
    width = Pt(element['width'])  # 使用点数
    height = Pt(element['height'])  # 使用点数
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = element['content']

    if 'font_name' in element and element['font_name']:
        run.font.name = element['font_name']
    if 'font_size' in element and element['font_size'] is not None:
        run.font.size = Pt(element['font_size'])  # 确保 font_size 使用点数
    if 'bold' in element:
        run.font.bold = element['bold']
    if 'italic' in element:
        run.font.italic = element['italic']
    if 'underline' in element:
        run.font.underline = element['underline']
    if 'color' in element and element['color']:
        run.font.color.rgb = RGBColor.from_string(element['color'].strip("#"))

def add_image(slide, element):
    image_path = element['path']
    left = Pt(element['left'])  # 使用点数
    top = Pt(element['top'])  # 使用点数
    width = Pt(element['width'])  # 使用点数
    height = Pt(element['height'])  # 使用点数
    slide.shapes.add_picture(image_path, left, top, width, height)


def json_to_ppt(json_path, pptx_path):
    with open(json_path, 'r') as f:
        data = json.load(f)

    ppt = Presentation()
    # 设置幻灯片尺寸
    ppt.slide_width = Pt(data["presentation"]["slide_size"]["width"])
    ppt.slide_height = Pt(data["presentation"]["slide_size"]["height"])

    for slide_data in data['presentation']['slides']:
        slide_layout = ppt.slide_layouts[5]  # Choosing a blank layout
        slide = ppt.slides.add_slide(slide_layout)

        # 按照 z_index 对元素进行排序
        sorted_elements = sorted(slide_data['elements'], key=lambda x: x.get('z_index', 0))

        for element in sorted_elements:
            if element['type'] == 'text':
                add_text(slide, element)
            elif element['type'] == 'image':
                add_image(slide, element)

    ppt.save(pptx_path)

# 路径可能需要根据您的实际文件位置进行调整
json_path = '/Users/birdmanoutman/PycharmProjects/cardesign-ppt-automaker/PPT_JSON_conversion/output_presentation.json'
pptx_path = 'recreated_presentation.pptx'
json_to_ppt(json_path, pptx_path)
