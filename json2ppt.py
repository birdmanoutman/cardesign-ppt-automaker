import json
import pptx

def add_elements(slide, elements_data):
    for element in elements_data:
        if element['type'] == 'text':
            left, top, width, height = element['left'], element['top'], element['width'], element['height']
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            paragraph = text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = element['text']
            if element['font_name']:
                run.font.name = element['font_name']
            if element['font_size']:
                run.font.size = pptx.util.Pt(element['font_size'])
            if element['space_after']:
                paragraph.space_after = element['space_after']

        elif element['type'] == 'image':
            image_path = element['path']
            left, top, width, height = element['left'], element['top'], element['width'], element['height']
            crop_left, crop_right, crop_top, crop_bottom = element['crop_left'], element['crop_right'], element['crop_top'], element['crop_bottom']
            picture = slide.shapes.add_picture(image_path, left, top, width, height)
            picture.crop_left = crop_left
            picture.crop_right = crop_right
            picture.crop_top = crop_top
            picture.crop_bottom = crop_bottom


def json_to_ppt(json_path):
    with open(json_path, "r") as f:
        data = json.load(f)
    ppt = pptx.Presentation()

    # Set slide width and height
    slide_width = pptx.util.Inches(data['slide_width'])
    slide_height = pptx.util.Inches(data['slide_height'])
    ppt.slide_layouts[0].slide_master.slide_size.width = slide_width
    ppt.slide_layouts[0].slide_master.slide_size.height = slide_height

    # Rest of the code remains the same


    slide_layouts = ppt.slide_layouts
    for slide_data in data['slides']:
        slide_layout_name = slide_data['slide_layout']
        slide_layout = [layout for layout in slide_layouts if layout.name == slide_layout_name][0]
        slide = ppt.slides.add_slide(slide_layout)
        elements_data = slide_data['elements']
        add_elements(slide, elements_data)
    return ppt

json_path = "example.json"
output_ppt = "recreated_example.pptx"

ppt = json_to_ppt(json_path)
ppt.save(output_ppt)
