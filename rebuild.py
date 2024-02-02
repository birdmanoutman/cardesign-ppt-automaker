import json
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import pptx

def set_slide_size(ppt, slide_size_data):
    ppt.slide_width = Inches(slide_size_data['width'])
    ppt.slide_height = Inches(slide_size_data['height'])

def find_layout_by_name(ppt, layout_name):
    for layout in ppt.slide_layouts:
        if layout.name == layout_name:
            return layout
    return None

def add_text(slide, text_data):
    for data in text_data:
        left, top, width, height = data['left'], data['top'], data['width'], data['height']
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        paragraph = text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = data['text']
        if data['font_name']:
            run.font.name = data['font_name']
        if data['font_size']:
            run.font.size = pptx.util.Pt(data['font_size'])
        if data['space_after']:
            run.space_after = data['space_after']

def add_images(slide, image_data):
    for data in image_data:
        image_path = data['path']
        left, top = data['left'], data['top']
        width, height = data['width'], data['height']
        crop_left, crop_right, crop_top, crop_bottom = data['crop_left'], data['crop_right'], data['crop_top'], data['crop_bottom']
        picture = slide.shapes.add_picture(image_path, left, top, width, height)
        picture.crop_left = crop_left
        picture.crop_right = crop_right
        picture.crop_top = crop_top
        picture.crop_bottom = crop_bottom


def add_elements(slide, elements_data):
    for element in elements_data:
        if element['type'] == 'text':
            left, top, width, height = element['left'], element['top'], element['width'], element['height']
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            paragraph = text_frame.add_paragraph()
            run = paragraph.add_run()
            run.text = element['text']
            if element['font_name']:
                run.font.name = element['font_name']
            if element['font_size']:
                run.font.size = pptx.util.Pt(element['font_size'])
            if element['space_after']:
                paragraph.space_after = element['space_after']

        elif element['type'] == 'image':
            image_path = element['path']
            left, top, width, height = element['left'], element['top'], element['width'], element['height']
            crop_left, crop_right, crop_top, crop_bottom = element['crop_left'], element['crop_right'], element['crop_top'], element['crop_bottom']
            picture = slide.shapes.add_picture(image_path, left, top, width, height)
            picture.crop_left = crop_left
            picture.crop_right = crop_right
            picture.crop_top = crop_top
            picture.crop_bottom = crop_bottom


def create_ppt_from_json(json_file, output_ppt):
    with open(json_file, "r") as f:
        presentation_data = json.load(f)

    ppt = Presentation()

    set_slide_size(ppt, presentation_data['slide_size'])

    for slide_data in presentation_data['slides']:
        layout_name = slide_data['layout']
        slide_layout = find_layout_by_name(ppt, layout_name)
        if slide_layout is None:
            slide_layout = ppt.slide_layouts[5]  # Use a blank slide layout as a fallback

        slide = ppt.slides.add_slide(slide_layout)

        add_text(slide, slide_data['text_data'])
        add_images(slide, slide_data['image_data'])

    ppt.save(output_ppt)


input_json = "example.json"
output_ppt = "recreated_example.pptx"

create_ppt_from_json(input_json, output_ppt)